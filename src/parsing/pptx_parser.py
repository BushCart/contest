from pptx import Presentation
from pathlib import Path


def parse_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    result = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            result.append({
                "text": "\n".join(texts),
                "source": Path(path).name,
                "page": i
            })

    return result
